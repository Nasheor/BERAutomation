#!/usr/bin/env python3
"""Add speaker notes to docs/environ_2026_presentation.pptx."""

from pptx import Presentation
from pptx.util import Pt

PPTX = r"C:\Users\nagar\PyCharmProjects\BERAutomation\docs\environ_2026_presentation.pptx"

# ── Speaker notes, one entry per slide ────────────────────────────────────────
NOTES = [

# ── Slide 1 — Title ──────────────────────────────────────────────────────────
"""\
Good morning — and thank you for the introduction. My name is Avinash Nagarajan, \
and I'm presenting work from the CIRCUS Interreg North-West Europe project at Munster \
Technological University.

What I want to show you today is a tool that can estimate the energy performance of any \
residential building across eight European countries — using nothing but a postal address, \
and delivering results in under thirty seconds. I'll walk you through the problem it \
addresses, how it works, and what it means in practice for communities trying to \
decarbonise their housing stock.\
""",

# ── Slide 2 — The Problem ─────────────────────────────────────────────────────
"""\
Let me set the scene. Buildings are responsible for roughly forty percent of total energy \
use in the EU — they are the single largest source of CO₂ in most North-West European \
countries. Improving the energy performance of existing housing stock is central to the \
EU's 2050 climate targets and, closer to home, to the ambitions of the CIRCUS project.

The primary tool for this work is the Energy Performance Certificate. These certificates — \
called a BER in Ireland, a DPE in France, an Energieausweis in Germany and Austria — tell \
you how energy-efficient a building is and guide retrofit decisions.

The problem is access. Getting an official EPC means paying a qualified assessor to visit \
the property, conduct a detailed survey, and process it through certification software. \
In Ireland that costs between €150 and €500 per house, and takes a specialist the better \
part of a day. For a community energy organisation trying to prioritise retrofit activities \
across an entire village or townland — usually with limited budgets and volunteer capacity — \
this is an enormous barrier. You simply cannot afford to pre-assess hundreds of homes \
before knowing which ones to target.\
""",

# ── Slide 3 — The Opportunity ─────────────────────────────────────────────────
"""\
But here is where it gets interesting. Satellite imagery now covers virtually every \
residential address in North-West Europe at a resolution that lets you measure a building's \
footprint. Google Street View provides near-complete photographic coverage of residential \
streets — effectively giving you an exterior survey of most properties without ever \
visiting them. Modern AI vision models can examine those photographs and make reliable \
inferences about when a building was constructed, what type it is, and what heating \
system it likely uses. And the energy calculation itself is not new — the HWB annual \
balance method has been used by practitioners for decades.

So we had all the ingredients: imagery, AI, and a well-established calculation method. \
The question was simply whether we could connect them into a fully automated pipeline — \
from postal code to BER rating, at scale, at near-zero marginal cost.\
""",

# ── Slide 4 — What the Tool Does ──────────────────────────────────────────────
"""\
The answer is yes — and this is what it looks like in practice. You open a browser, \
type in a postal code or street address, select the country, and click run. \
That is the entire user interaction required.

Within about thirty seconds, the tool returns a full building energy profile. You get \
the BER band on the Irish scale — which we use as a consistent cross-border reference \
for the CIRCUS programme — alongside the native EPC format that residents in each country \
would actually recognise on their own national certificate. You get primary energy in \
kilowatt-hours per square metre per year, annual CO₂ emissions, a breakdown of where \
energy is being lost, and the building's physical characteristics as estimated by the AI. \
And you get a retrofit comparison: what would the rating become if you added insulation \
or switched to a heat pump?

The tool covers all eight partner countries of the CIRCUS Interreg NWE programme — Ireland, \
France, Germany, Belgium, the Netherlands, Luxembourg, Switzerland, and Austria. \
No specialist knowledge is required to operate it. A community energy officer or local \
authority staff member can run it directly.\
""",

# ── Slide 5 — The Pipeline ────────────────────────────────────────────────────
"""\
Let me walk through the pipeline, because the engineering decisions are worth understanding.

Phase one is geocoding — the address is converted to GPS coordinates using the Google \
Geocoding API, with a country filter applied to ensure we stay within the right borders.

Phase two fetches imagery. A satellite image at zoom level twenty gives us a top-down view \
of the building at roughly nine centimetres per pixel. For Street View, rather than using \
a fixed compass direction, the tool queries the Street View metadata API to find where the \
nearest camera physically was, computes the geodesic bearing from that camera to the \
building, and then fetches four images at ninety-degree intervals around that bearing. \
The building is centred in every photograph, and we capture the full exterior — front, \
right, rear, and left.

Phase three sends all four images simultaneously to Claude Sonnet, Anthropic's AI vision \
model, along with a detailed country-specific prompt. Claude returns its assessment of \
building type, construction era, storeys, heating system, and a confidence score.

Phase four extracts the building's footprint from the satellite image — also using Claude \
Vision, with the map scale in metres per pixel injected into the prompt so it can reason \
in real-world dimensions. An independent OpenCV algorithm runs in parallel as a \
cross-validation check. If both methods agree within thirty percent, confidence is boosted; \
if they disagree, the AI result is trusted; if the AI fails entirely, OpenCV acts as fallback.

Phase five runs the energy calculation — which I will explain on the next slide.\
""",

# ── Slide 6 — AI Building Analysis ───────────────────────────────────────────
"""\
The building analysis phase is where the AI does the heavy lifting, and a few design \
decisions here are particularly important.

Sending all four street view images in one request means the model can cross-reference \
them. An oil tank visible only from the rear, a heat pump unit on the side wall, a shared \
party wall only apparent from the left — the model picks all of these up when it has a \
complete exterior survey, rather than making inferences from a single front-facing \
photograph.

We have also embedded country-specific architectural knowledge directly into the prompts. \
The visual signatures of construction eras differ significantly between countries — a 1980s \
German building with heavy external insulation render looks nothing like a 1980s Irish \
cavity-wall semi-detached. Rural Irish and French properties are far more likely to have \
oil storage tanks. District heating is common in Luxembourg and the Netherlands but rare \
in rural Ireland. This localisation meaningfully improves classification accuracy across \
all eight countries.

The confidence gating is a critical safeguard. The model returns a confidence score with \
every analysis. If that score falls below 0.4 — the building is obscured by vegetation, \
or there is simply no Street View coverage at that address — the tool does not use the \
uncertain classification. It applies conservative defaults and displays a clear warning \
to the user. We would rather give an honest, conservative default than propagate a \
low-quality AI guess through to the final energy figure.\
""",

# ── Slide 7 — The Energy Calculation ─────────────────────────────────────────
"""\
The energy calculation uses the HWB annual heating balance — Heizwärmebedarf in German, \
meaning heating energy requirement. This is a steady-state annual method based on \
ISO 13790, the same standard that underpins professional EPC calculation tools across \
Europe. It was ported cell by cell from an Excel workbook developed by Benjamin Kaiser \
at MTU in January 2025.

The balance is straightforward: heating demand equals transmission losses through the \
building fabric plus ventilation losses from air exchange, minus usable solar gains \
through glazing and internal heat from occupants and appliances. That heating demand is \
divided by the heating system efficiency — or the seasonal coefficient of performance for \
heat pumps, which can exceed 1.0 because heat pumps extract free environmental energy — \
to give final energy consumption. Final energy is then multiplied by a country-specific \
Primary Energy Factor to get primary energy, which determines the EPC band.

Every country-specific input comes from authoritative sources: heating degree days from \
degreedays.net, solar irradiance from the Passive House Planning Package climate database, \
primary energy factors from each country's national EPBD transposition, and grid carbon \
intensity from national grid operators. The same physical building in different countries \
produces a genuinely different calculation — not just a different label.\
""",

# ── Slide 8 — CO₂ Across Countries ───────────────────────────────────────────
"""\
This is the slide I want you to pause on, because I think it carries the most important \
policy implication of everything I am showing you today.

We know that heat pumps are central to retrofit strategies across Europe. But what we \
often miss is that installing a heat pump produces very different CO₂ outcomes depending \
on where you are — because the carbon intensity of the electricity grid varies enormously \
between countries.

Switzerland's grid runs on hydropower and nuclear, producing just twenty-nine grams of \
CO₂ per kilowatt-hour. France's nuclear-heavy grid sits at fifty-two grams. At the other \
end, Germany's coal and gas mix emits three hundred and eighty-five grams per kilowatt-hour \
— more than thirteen times Switzerland's figure.

What does this mean in practice? Install an air-source heat pump with a seasonal COP of \
3.5 in a French home, and you cut CO₂ by roughly eighty-nine percent compared to a gas \
boiler. The same heat pump in Germany cuts CO₂ by about forty-nine percent — still \
significant, but not transformational, and only until Germany's grid catches up.

For CIRCUS partners advising communities on where to focus retrofit resources, this \
distinction is critical. In France and Switzerland, pushing hard on heat pump uptake makes \
sense right now. In Germany and the Netherlands, an insulation-first strategy may deliver \
more immediate CO₂ benefit, with heat pumps becoming progressively more powerful as the \
grid decarbonises. The tool makes this visible, by country, for every assessment it runs.\
""",

# ── Slide 9 — Retrofit Analysis ───────────────────────────────────────────────
"""\
The retrofit analysis capability is what makes the tool actionable rather than just \
descriptive.

For any building assessed through the pipeline, you can specify a retrofit package — \
wall and roof insulation thickness, a target window U-value, and whether to switch the \
heating system — and the tool recalculates the BER side by side with the current rating.

The example here is a pre-1980 detached Irish house on a gas boiler — one of the most \
common and most energy-intensive typologies in rural Ireland. Its baseline sits at C2, \
at 177.6 kilowatt-hours per square metre per year. Add wall and roof insulation plus \
triple-glazed windows, and it moves to B3 — a twenty-seven percent reduction in CO₂. \
Add a heat pump on top of that, and it reaches A2 — an eighty-six percent cut in CO₂ \
compared to the unimproved baseline.

For a community energy officer running a public meeting, this tool lets them pull up \
any address, show the estimated current rating, and immediately demonstrate what is \
achievable with a standard deep retrofit package. It makes the pathway to an A-rated home \
tangible and evidence-based for homeowners who may never have seen their building's \
energy profile before. That is a powerful engagement tool.\
""",

# ── Slide 10 — Accuracy, Limitations & What's Next ──────────────────────────
"""\
I want to be transparent about what the tool is and what it is not.

For a building with good Street View coverage and a clearly visible facade, we expect \
agreement within one to two BER bands of an official assessment — which is a useful and \
actionable level of accuracy for screening purposes. That widens to two to three bands \
for partially obscured buildings, and three to four bands for rural properties with no \
Street View coverage at all.

This is intentionally a screening and prioritisation tool. It should not be used as a \
substitute for an official EPC where one is legally required — property transactions, \
building regulations compliance, green finance products. For those purposes, a qualified \
assessor must conduct a formal assessment. But for identifying where the worst-performing \
homes are in a community, preparing evidence for funding applications, or showing \
homeowners their retrofit pathway — this is exactly the right level of detail.

On next steps: we are currently validating the tool against the SEAI BER certificate \
database, which contains over four hundred thousand rated Irish properties. That will allow \
us to quantify systematic biases and calibrate confidence across different building \
typologies. We are also developing batch processing so community organisations can upload \
a list of addresses and receive assessments in bulk.

Most importantly, we are actively looking for collaboration — with CIRCUS partners across \
all eight countries, with SEAI, with local authorities, and with researchers in this room. \
If any of that connects with your work, I would very much welcome a conversation \
afterwards. Thank you.\
""",

]

# ── Write notes into the presentation ─────────────────────────────────────────
prs = Presentation(PPTX)

for slide, note_text in zip(prs.slides, NOTES):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame

    # Clear any existing notes
    for para in tf.paragraphs[1:]:
        p_elem = para._p
        p_elem.getparent().remove(p_elem)

    tf.paragraphs[0].text = ""   # clear first paragraph text

    # Split note into paragraphs at blank lines
    blocks = note_text.split("\n\n")
    first = True
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        if first:
            p = tf.paragraphs[0]
            p.text = block
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(12)
            first = False
        else:
            p = tf.add_paragraph()
            p.text = block
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(12)

prs.save(PPTX)
print(f"Notes added to all {len(NOTES)} slides.")
print(f"Saved: {PPTX}")
