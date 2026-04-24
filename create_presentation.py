#!/usr/bin/env python3
"""Generate docs/environ_2026_presentation.pptx for Environ 2026 conference."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE = r"C:\Users\nagar\PyCharmProjects\BERAutomation"
LOGO = BASE + r"\images\circus_logo-presentation.png"
OUT  = BASE + r"\docs\environ_2026_presentation.pptx"

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
GREY   = RGBColor(0x55, 0x55, 0x55)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
LBLUE  = RGBColor(0xA8, 0xC8, 0xE8)
LGREY  = RGBColor(0xE8, 0xEC, 0xF2)

# Logo aspect ratio: 734 x 340 px → W/H = 2.159
LOGO_ASPECT = 734 / 340

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]   # blank layout

# ── Utility helpers ───────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_logo(slide, width=Inches(1.8)):
    height = width / LOGO_ASPECT
    slide.shapes.add_picture(
        LOGO,
        SW - width - Inches(0.18),
        SH - height - Inches(0.12),
        width,
        height,
    )


def _tf(slide, l, t, w, h):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    return tf


def run(para, text, size, bold=False, italic=False, color=DARK):
    r = para.add_run()
    r.text = text
    r.font.size    = Pt(size)
    r.font.bold    = bold
    r.font.italic  = italic
    r.font.color.rgb = color
    return r


def content_header(slide, title, num):
    """Navy bar with white title and slide counter."""
    bar_h = Inches(1.05)
    add_rect(slide, 0, 0, SW, bar_h, NAVY)
    add_rect(slide, 0, bar_h, SW, Inches(0.045), GREEN)   # green accent line

    tf = _tf(slide, Inches(0.45), Inches(0.13), SW - Inches(2.8), bar_h - Inches(0.13))
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run(p, title, 28, bold=True, color=WHITE)

    tf2 = _tf(slide, SW - Inches(1.1), Inches(0.07), Inches(0.95), Inches(0.4))
    p2  = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run(p2, f"{num} / 10", 12, color=LBLUE)
    return bar_h


def bullets(slide, items, top, left=Inches(0.5), width=None,
            main_sz=19, sub_sz=17):
    """items = list of (text, level) where level 0=main, 1=sub."""
    if width is None:
        width = SW - Inches(0.85)
    tf = _tf(slide, left, top, width, SH - top - Inches(0.95))
    first = True
    for text, level in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6 if level == 0 else 2)
        prefix = "•  " if level == 0 else "   –  "
        sz     = main_sz   if level == 0 else sub_sz
        col    = DARK      if level == 0 else GREY
        run(p, prefix + text, sz, color=col)
    return tf


def make_table(slide, headers, rows, l, t, w, h,
               hdr_color=NAVY, hdr_sz=14, cell_sz=13):
    ncols  = len(headers)
    nrows  = len(rows) + 1
    tbl    = slide.shapes.add_table(nrows, ncols, l, t, w, h).table

    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hdr
        r.font.bold = True
        r.font.size = Pt(hdr_sz)
        r.font.color.rgb = WHITE

    for i, row in enumerate(rows):
        bg = WHITE if i % 2 == 0 else LGREY
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            r.font.size = Pt(cell_sz)
            r.font.color.rgb = DARK

    return tbl


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_rect(slide, 0, 0, SW, SH, NAVY)                        # full navy background
add_rect(slide, 0, 0, SW, Inches(0.07), GREEN)              # top accent

# Main title
tf = _tf(slide, Inches(0.7), Inches(1.1), SW - Inches(4.3), Inches(2.8))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run(p, "BER Automation:", 44, bold=True, color=WHITE)
p2 = tf.add_paragraph()
p2.space_before = Pt(10)
p2.alignment = PP_ALIGN.LEFT
run(p2, "Estimating Building Energy Ratings\nat Scale Using AI", 32, bold=True, color=LBLUE)

# Divider
add_rect(slide, Inches(0.7), Inches(3.95), Inches(5.5), Inches(0.045), GREEN)

# Details block
tf2 = _tf(slide, Inches(0.7), Inches(4.1), SW - Inches(4.3), Inches(2.5))
details = [
    ("Environ 2026",                                                          22, True,  WHITE),
    ("Avinash Nagarajan  ·  Munster Technological University",           16, False, LBLUE),
    ("CIRCUS — Connecting Initiatives for Rural Communities,",           15, False, LBLUE),
    ("Upscaling their Sustainable Energy  ·  Interreg NWE",             15, False, LBLUE),
]
first = True
for text, sz, bold, col in details:
    if first:
        p = tf2.paragraphs[0]; first = False
    else:
        p = tf2.add_paragraph()
    p.space_before = Pt(8 if text.startswith("Environ") else 4)
    p.alignment = PP_ALIGN.LEFT
    run(p, text, sz, bold=bold, color=col)

# Logo — larger on title slide
logo_w = Inches(3.2)
logo_h = logo_w / LOGO_ASPECT
slide.shapes.add_picture(LOGO,
    SW - logo_w - Inches(0.5),
    SH - logo_h - Inches(0.35),
    logo_w, logo_h)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
content_header(slide, "The Problem", 2)

items = [
    ("Buildings account for 40% of total energy consumption in the EU", 0),
    ("The largest single source of CO₂ in most North-West European countries", 1),
    ("Energy Performance Certificates (EPCs) are the primary tool for communicating a building's energy status", 0),
    ("BER (Ireland) · DPE (France) · Energieausweis (Germany/Austria) · Energielabel (Netherlands)…", 1),
    ("Commissioning an official EPC requires:", 0),
    ("A qualified assessor visiting the property for 1–3 hours", 1),
    ("Specialist certification software", 1),
    ("€150–€500 per property", 1),
    ("This creates a major barrier for:", 0),
    ("Community energy organisations prioritising retrofit across hundreds of homes", 1),
    ("Local authorities conducting large-scale housing stock assessments", 1),
    ("CIRCUS Interreg NWE partners comparing energy performance across 8 countries", 1),
]
bullets(slide, items, top=Inches(1.2))
add_logo(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Opportunity
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
content_header(slide, "The Opportunity", 3)

items = [
    ("Satellite imagery covers virtually every residential address in NWE — at ~9 cm/pixel resolution", 0),
    ("Google Street View provides photographic exterior surveys of most properties", 0),
    ("AI vision models can classify building features from photographs with remarkable reliability", 0),
    ("The energy physics is well-established — the HWB calculation method has been used by practitioners for decades", 0),
    ("All the inputs needed for an energy rating exist in publicly accessible data", 0),
]
bullets(slide, items, top=Inches(1.25), main_sz=21)

# Bold question
tf = _tf(slide, Inches(0.5), Inches(5.0), SW - Inches(1.0), Inches(1.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run(p,
    "→  Can we automate the entire workflow — from postal address to BER rating —"
    " using only publicly available data and AI vision?",
    22, bold=True, color=NAVY)

add_logo(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What the Tool Does
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
content_header(slide, "What the Tool Does", 4)

# Left: input / output
tf_l = _tf(slide, Inches(0.45), Inches(1.2), Inches(7.6), Inches(5.8))
p = tf_l.paragraphs[0]
run(p, "Input:", 20, bold=True, color=NAVY)
for item in ["A postal code or street address  +  country selection", "Nothing else."]:
    p2 = tf_l.add_paragraph()
    p2.space_before = Pt(4)
    run(p2, "•  " + item, 18, color=DARK)

p_out = tf_l.add_paragraph()
p_out.space_before = Pt(16)
run(p_out, "Output in under 30 seconds:", 20, bold=True, color=NAVY)

out_items = [
    "BER band (Irish A1–G as cross-border reference scale)",
    "Native country EPC band (DPE / Energieausweis / Energielabel / GEAK…)",
    "Primary energy: kWh/m²/year  ·  CO₂ emissions: kg/year",
    "Full energy breakdown: transmission, ventilation, solar & internal gains",
    "Building characteristics: dimensions, floor area, type, era, heating system",
    "Retrofit comparison: rating after insulation + heating system upgrade",
]
for item in out_items:
    p2 = tf_l.add_paragraph()
    p2.space_before = Pt(5)
    run(p2, "•  " + item, 17, color=DARK)

# Right: supported countries box
add_rect(slide, Inches(8.3), Inches(1.2), Inches(4.7), Inches(5.7), LGREY)
tf_r = _tf(slide, Inches(8.5), Inches(1.3), Inches(4.3), Inches(5.5))
p = tf_r.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run(p, "8 Supported Countries", 19, bold=True, color=NAVY)

countries = [
    "IE  Ireland (BER A1–G)",
    "FR  France (DPE A–G)",
    "DE  Germany (Energieausweis A+–H)",
    "BE  Belgium (EPC A+–F)",
    "NL  Netherlands (Energielabel A++++–G)",
    "LU  Luxembourg (Energiepass A+–G)",
    "CH  Switzerland (GEAK A–G)",
    "AT  Austria (Energieausweis A++–G+)",
]
for c in countries:
    p2 = tf_r.add_paragraph()
    p2.space_before = Pt(9)
    p2.alignment = PP_ALIGN.LEFT
    run(p2, "•  " + c, 16, color=DARK)

add_logo(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Pipeline
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
content_header(slide, "The Automated Pipeline", 5)

# Five phase boxes
phases = [
    ("1  GEOCODING",
     "Address → GPS\n(Google Geocoding API)",
     NAVY),
    ("2  IMAGE\n    FETCHING",
     "Satellite (zoom 20)\n+ 4× Street View\n0°/90°/180°/270°",
     RGBColor(0x1A, 0x5C, 0x8A)),
    ("3  AI BUILDING\n    ANALYSIS",
     "Claude Vision →\ntype, era, storeys,\nheating system",
     RGBColor(0x2E, 0x60, 0x7B)),
    ("4  FOOTPRINT\n    EXTRACTION",
     "Claude Vision (primary)\n+ OpenCV\ncross-validation",
     RGBColor(0x1A, 0x6B, 0x5A)),
    ("5  HWB\n    CALCULATION",
     "Annual energy balance\n→ kWh/m²/yr\n→ BER band",
     GREEN),
]

box_w = Inches(2.28)
box_h = Inches(1.72)
gap   = Inches(0.22)
total = len(phases) * box_w + (len(phases) - 1) * gap
start = (SW - total) / 2
top_b = Inches(1.35)

for i, (label, detail, color) in enumerate(phases):
    l = start + i * (box_w + gap)
    add_rect(slide, l, top_b, box_w, box_h, color)
    tf = _tf(slide, l + Inches(0.1), top_b + Inches(0.1),
             box_w - Inches(0.2), box_h - Inches(0.2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, label, 13, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    p2.alignment = PP_ALIGN.CENTER
    run(p2, detail, 11, color=LBLUE)
    # Arrow
    if i < len(phases) - 1:
        ax = l + box_w
        ay = top_b + box_h / 2 - Inches(0.04)
        add_rect(slide, ax, ay, gap, Inches(0.055), GREY)

# Labels
tf_in = _tf(slide, start, Inches(1.08), Inches(5.0), Inches(0.3))
run(tf_in.paragraphs[0], "Input: address + country", 14, italic=True, color=GREY)

last_l = start + (len(phases) - 1) * (box_w + gap)
tf_out = _tf(slide, last_l - Inches(1.5), top_b + box_h + Inches(0.08), Inches(4.0), Inches(0.3))
p_out = tf_out.paragraphs[0]
p_out.alignment = PP_ALIGN.RIGHT
run(p_out, "Output: BER + native EPC + CO₂", 14, italic=True, color=GREY)

# Key callouts
key_items = [
    ("No site visit.  No specialist.  Near-zero marginal cost per assessment.", 0),
    ("Confidence gating: if AI confidence < 0.4 the tool uses safe defaults — it never propagates a low-quality guess", 0),
    ("Graceful degradation at every phase: each step has a fallback so one failure does not crash the pipeline", 0),
]
bullets(slide, key_items, top=Inches(3.3), main_sz=17)
add_logo(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AI Building Analysis
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
content_header(slide, "AI Building Analysis", 6)

left_items = [
    ("Model: Claude Sonnet (Anthropic) — multimodal vision + language", 0),
    ("All 4 street view images sent in one request — model cross-references all angles", 0),
    ("Features visible from only one angle are captured: oil tanks, heat pump units, party walls, flues", 1),
    ("Country-specific architectural knowledge embedded in prompts:", 0),
    ("Ireland / France: look for external oil tanks in rural areas", 1),
    ("Luxembourg / Netherlands: district heating dominant in urban areas", 1),
    ("Germany: heavy external insulation render common post-1980", 1),
    ("Switzerland: triple glazing + heat pumps more common post-2010 (Minergie)", 1),
]
bullets(slide, left_items, top=Inches(1.2), width=Inches(6.8), main_sz=18, sub_sz=16)

# Right: AI output table
tbl_l = Inches(7.2)
tbl_w = SW - tbl_l - Inches(0.25)
tbl = make_table(
    slide,
    ["Parameter", "Examples"],
    [
        ["Building type",    "Detached / Semi-D / Terraced"],
        ["Construction era", "Before 1980 → After 2010"],
        ["Storeys",          "1 / 2 / 3 / 4"],
        ["Heating system",   "Oil, gas, heat pump, biomass, district"],
        ["Confidence",       "0.0 – 1.0"],
        ["Reasoning",        "Plain-language explanation"],
    ],
    tbl_l, Inches(1.2), tbl_w, Inches(3.0),
    hdr_sz=15, cell_sz=13,
)
tbl.columns[0].width = Inches(2.15)
tbl.columns[1].width = tbl_w - Inches(2.15)

# Confidence note
tf = _tf(slide, Inches(0.45), Inches(5.55), SW - Inches(1.0), Inches(0.85))
p = tf.paragraphs[0]
run(p, "Confidence gating: ", 17, bold=True, color=NAVY)
run(p, "if confidence < 0.4 (building obscured or no Street View coverage), the tool applies safe defaults "
       "and shows the user a clear warning. It does not propagate an uncertain AI guess into the energy calculation.",
    17, color=DARK)

add_logo(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — The Energy Calculation
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
content_header(slide, "The Energy Calculation (HWB Method)", 7)

left_items = [
    ("HWB (Heizwärmebedarf) annual heating balance", 0),
    ("ISO 13790 steady-state basis — same standard used by professional EPC tools across Europe", 1),
    ("Ported cell-by-cell from MTU Excel tool (Benjamin Kaiser, January 2025)", 1),
    ("Every country-specific input sourced from authoritative data:", 0),
    ("Heating Degree Days — degreedays.net 2022 data", 1),
    ("Solar irradiance — PHPP climate database", 1),
    ("Primary Energy Factor — each country's national EPBD transposition", 1),
    ("Grid CO₂ intensity — national grid operators (SEAI, RTE, UBA, ELIA, RVO…)", 1),
]
bullets(slide, left_items, top=Inches(1.2), width=Inches(6.8), main_sz=18, sub_sz=16)

# Right: formula box
add_rect(slide, Inches(7.2), Inches(1.2), Inches(5.85), Inches(5.2), LGREY)
tf = _tf(slide, Inches(7.35), Inches(1.3), Inches(5.6), Inches(5.0))

formula_lines = [
    ("Annual energy balance:", 16, True,  NAVY),
    ("", 8, False, DARK),
    ("Q_heating  =  Transmission losses (fabric)", 13, False, DARK),
    ("          +  Ventilation losses (air exchange)", 13, False, DARK),
    ("          −  Solar gains (passive solar through windows)", 13, False, DARK),
    ("          −  Internal gains (occupants + appliances)", 13, False, DARK),
    ("", 8, False, DARK),
    ("Final energy  =  Q_heating  ÷  System efficiency", 13, False, DARK),
    ("  (SCOP for heat pumps — can exceed 1.0)", 12, True, GREY),
    ("", 8, False, DARK),
    ("Primary energy  =  Final energy  ×  PEF  (country-specific)", 13, False, DARK),
    ("", 8, False, DARK),
    ("CO₂  =  Final energy  ×  Grid emission factor", 13, False, DARK),
    ("", 8, False, DARK),
    ("BER band  =  lookup( primary energy / floor area )", 14, True, NAVY),
]
first = True
for text, sz, bold, col in formula_lines:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run(p, text, sz, bold=bold, color=col)

add_logo(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CO₂ Across Countries
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
content_header(slide, "CO₂ Across Countries: A Policy Insight", 8)

tbl_w = Inches(7.3)
tbl = make_table(
    slide,
    ["Country", "Grid CO₂ (g/kWh)", "Heat pump CO₂ reduction vs gas"],
    [
        ["Switzerland",  "29",   "~93%"],
        ["France",       "52",   "~89%"],
        ["Belgium",      "163",  "~72%"],
        ["Austria",      "156",  "~73%"],
        ["Luxembourg",   "197",  "~66%"],
        ["Ireland",      "210",  "~64%"],
        ["Netherlands",  "290",  "~51%"],
        ["Germany",      "385",  "~49%"],
    ],
    Inches(0.45), Inches(1.25), tbl_w, Inches(4.0),
    hdr_sz=15, cell_sz=14,
)
tbl.columns[0].width = Inches(2.4)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = tbl_w - Inches(4.4)

# Highlight best (Switzerland, row 1) and worst (Germany, row 8)
for ci in range(3):
    tbl.cell(1, ci).fill.solid()
    tbl.cell(1, ci).fill.fore_color.rgb = RGBColor(0xC8, 0xE6, 0xC9)   # light green
    tbl.cell(8, ci).fill.solid()
    tbl.cell(8, ci).fill.fore_color.rgb = RGBColor(0xFF, 0xCC, 0xBC)   # light amber

# Right: insight text
tf_r = _tf(slide, Inches(7.95), Inches(1.25), Inches(5.1), Inches(5.8))
p = tf_r.paragraphs[0]
run(p, "Same building.  Same heat pump.  Same SCOP of 3.5.", 19, bold=True, color=NAVY)

insights = [
    ("", 8, False, DARK),
    ("In France: a heat pump is near-zero-carbon heating today.", 17, False, DARK),
    ("", 5, False, DARK),
    ("In Germany: a heat pump halves CO₂ but is not yet transformational — grid decarbonisation must follow.", 17, False, DARK),
    ("", 5, False, DARK),
    ("For CIRCUS partners advising communities on retrofit priority, this distinction is critical.", 17, True, NAVY),
    ("", 5, False, DARK),
    ("Insulation-first strategies may be more immediately effective in high-carbon-grid countries.", 16, False, GREY),
]
for text, sz, bold, col in insights:
    p2 = tf_r.add_paragraph()
    p2.space_before = Pt(3)
    run(p2, text, sz, bold=bold, color=col)

add_logo(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Retrofit Analysis
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
content_header(slide, "Retrofit Analysis", 9)

# Sub-heading
tf = _tf(slide, Inches(0.45), Inches(1.2), SW - Inches(0.9), Inches(0.42))
p = tf.paragraphs[0]
run(p, "Example: Pre-1980 Irish detached house on gas boiler — a very common typology in rural Ireland",
    17, italic=True, color=GREY)

tbl_w = SW - Inches(0.9)
tbl = make_table(
    slide,
    ["Scenario", "Primary energy", "BER band", "CO₂"],
    [
        ["Baseline (pre-1980, gas boiler)",
         "177.6 kWh/m²/yr", "C2", "31.3 kg/m²/yr"],
        ["+ Wall insulation (12 cm)  +  Roof insulation (20 cm)  +  Triple glazing",
         "~130 kWh/m²/yr", "B3", "~22.9 kg/m²/yr"],
        ["+ Switch to air-source heat pump",
         "~50 kWh/m²/yr", "A2", "~4.3 kg/m²/yr"],
        ["Total improvement",
         "−72%", "C2 → A2", "−86% CO₂"],
    ],
    Inches(0.45), Inches(1.72), tbl_w, Inches(2.45),
    hdr_sz=15, cell_sz=14,
)
tbl.columns[0].width = Inches(5.5)
tbl.columns[1].width = Inches(2.35)
tbl.columns[2].width = Inches(1.5)
tbl.columns[3].width = tbl_w - Inches(9.35)

# Highlight total row (row index 4)
for ci in range(4):
    cell = tbl.cell(4, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for r in cell.text_frame.paragraphs[0].runs:
        r.font.color.rgb = WHITE
        r.font.bold = True

key_items = [
    ("Built-in before/after comparison — specify insulation thickness, window U-value, heating system swap", 0),
    ("Identifies which homes benefit most from insulation vs heating system change", 0),
    ("Makes the retrofit pathway tangible for homeowners and community energy organisations", 0),
    ("Directly supports national grant scheme (SEAI / equivalent) conversations", 0),
]
bullets(slide, key_items, top=Inches(4.38), main_sz=18)
add_logo(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Accuracy, Limitations & What's Next
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
content_header(slide, "Accuracy, Limitations & What’s Next", 10)

# Accuracy table (left)
tbl = make_table(
    slide,
    ["Scenario", "Expected accuracy"],
    [
        ["Clear building, good Street View coverage",  "±1–2 BER bands"],
        ["Obscured or limited Street View",            "±2–3 BER bands"],
        ["No Street View coverage (rural)",            "±3–4 BER bands"],
        ["Apartments / multi-unit blocks",             "Not recommended"],
    ],
    Inches(0.45), Inches(1.2), Inches(6.5), Inches(2.05),
    hdr_sz=15, cell_sz=14,
)
tbl.columns[0].width = Inches(4.5)
tbl.columns[1].width = Inches(2.0)

# Limitations note
tf = _tf(slide, Inches(0.45), Inches(3.38), Inches(6.5), Inches(0.45))
run(tf.paragraphs[0],
    "Screening tool — not a substitute for official EPC certification",
    17, bold=True, color=NAVY)

lim_items = [
    ("U-values from Austrian OIB standards — broadly comparable, not Irish Part L exact", 0),
    ("Electricity CO₂ factors from 2022 — grids are decarbonising rapidly", 0),
    ("Validation vs SEAI BER database (400k+ records) — ongoing", 0),
]
bullets(slide, lim_items, top=Inches(3.9), width=Inches(6.5), main_sz=17)

# Right: next steps box
add_rect(slide, Inches(7.15), Inches(1.2), Inches(5.88), Inches(5.2), LGREY)
tf_r = _tf(slide, Inches(7.35), Inches(1.3), Inches(5.55), Inches(4.9))
p = tf_r.paragraphs[0]
run(p, "What’s Next", 22, bold=True, color=NAVY)

next_steps = [
    "Validation against SEAI BER database",
    "Batch processing for estate-level assessment",
    "Regional climate data (county-level HDD for Ireland)",
    "DEAP method mode (official Irish methodology)",
    "PDF BER-style report generation",
]
for step in next_steps:
    p2 = tf_r.add_paragraph()
    p2.space_before = Pt(9)
    run(p2, "→  " + step, 16, color=DARK)

# Collaboration
p_div = tf_r.add_paragraph()
p_div.space_before = Pt(18)
run(p_div, "Open to collaboration:", 18, bold=True, color=NAVY)
for line in ["CIRCUS Interreg NWE partners", "SEAI & local authorities", "Environ researchers"]:
    p2 = tf_r.add_paragraph()
    p2.space_before = Pt(6)
    run(p2, "•  " + line, 16, color=DARK)

# Thank you
tf_ty = _tf(slide, Inches(0.45), Inches(5.95), Inches(6.5), Inches(0.85))
p = tf_ty.paragraphs[0]
run(p, "Thank you  —  Questions welcome", 22, bold=True, color=NAVY)

add_logo(slide)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
