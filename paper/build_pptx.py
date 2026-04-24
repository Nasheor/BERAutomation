"""
Build BER Automation presentation as a PowerPoint .pptx file.
Mirrors the 5-slide HTML deck.

Run:  python paper/build_pptx.py
Output: paper/BER_Automation_CIRCUS.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy
import os

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
TEAL   = RGBColor(0x00, 0xB4, 0xC6)
AMBER  = RGBColor(0xF4, 0xA4, 0x27)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
DARK   = RGBColor(0x1F, 0x2A, 0x37)

RED_LIGHT   = RGBColor(0xFE, 0xF2, 0xF2)
RED_BORDER  = RGBColor(0xFE, 0xCA, 0xCA)
RED_DARK    = RGBColor(0xB9, 0x1C, 0x1C)
GREEN_LIGHT = RGBColor(0xF0, 0xFD, 0xF4)
GREEN_DARK  = RGBColor(0x06, 0x5F, 0x46)
GREEN_MID   = RGBColor(0x16, 0xA3, 0x4A)
GREEN_BORDER= RGBColor(0xBB, 0xF7, 0xD0)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LIGHT= RGBColor(0xF5, 0xF3, 0xFF)
PURPLE_DARK = RGBColor(0x6B, 0x21, 0xA8)
BLUE        = RGBColor(0x23, 0x63, 0xEB)
BLUE_LIGHT  = RGBColor(0xEF, 0xF6, 0xFF)
BLUE_DARK   = RGBColor(0x1E, 0x40, 0xAF)
AMBER_LIGHT = RGBColor(0xFF, 0xFB, 0xEB)
AMBER_DARK  = RGBColor(0x92, 0x40, 0x0E)
ORANGE      = RGBColor(0xE9, 0x73, 0x16)
SLATE       = RGBColor(0x47, 0x55, 0x69)
SLATE_LIGHT = RGBColor(0xF0, 0xF4, 0xF8)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.width = line_width
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
                wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb


def add_label_value(slide, label, value, x, y, w=3.0, label_size=9, value_size=13):
    """Small metric block: label above, value below."""
    add_textbox(slide, label, x, y,       w, 0.28, size=label_size, color=GRAY,     align=PP_ALIGN.LEFT)
    add_textbox(slide, value, x, y+0.25,  w, 0.38, size=value_size, bold=True,  color=DARK, align=PP_ALIGN.LEFT)


def header(slide, title, badge_text=""):
    """Navy header bar with optional badge."""
    add_rect(slide, 0, 0, 13.33, 0.72, fill=NAVY)
    if badge_text:
        add_rect(slide, 0.22, 0.13, 0.52, 0.46, fill=TEAL)
        add_textbox(slide, badge_text, 0.22, 0.13, 0.52, 0.46,
                    size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, 0.85, 0.13, 12.0, 0.46,
                    size=18, bold=True, color=WHITE)
    else:
        add_textbox(slide, title, 0.3, 0.13, 12.7, 0.46,
                    size=18, bold=True, color=WHITE)


def footer(slide, left_text, right_text):
    """Light footer bar."""
    add_rect(slide, 0, 7.18, 13.33, 0.32, fill=LIGHT)
    add_textbox(slide, left_text,  0.22, 7.19, 6.0,  0.28, size=9, color=GRAY)
    add_textbox(slide, right_text, 7.0,  7.19, 6.0,  0.28, size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def colored_box_with_text(slide, x, y, w, h, bg, border, lines, line_sizes, line_colors, line_bolds=None, rx=0):
    """Draw a filled rectangle then overlay multiple text lines."""
    add_rect(slide, x, y, w, h, fill=bg, line_color=border, line_width=Pt(0.75))
    if line_bolds is None:
        line_bolds = [False] * len(lines)
    ty = y + 0.08
    for i, (txt, sz, col, bld) in enumerate(zip(lines, line_sizes, line_colors, line_bolds)):
        add_textbox(slide, txt, x + 0.1, ty, w - 0.18, 0.3,
                    size=sz, color=col, bold=bld)
        ty += 0.26


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════

def slide1(prs):
    sl = blank_slide(prs)

    # Full navy background
    add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)

    # Teal super-label
    add_textbox(sl, "CIRCUS PROJECT — TOOL DEMONSTRATION",
                0, 1.05, 13.33, 0.4, size=10, bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)

    # Main title
    add_textbox(sl, "Building Energy Rating\nAutomation Tool",
                0.5, 1.5, 12.33, 1.4, size=40, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(sl, "From postcode to energy rating — in under 30 seconds",
                0, 2.95, 13.33, 0.5, size=16, color=GRAY,
                align=PP_ALIGN.CENTER)

    # Before/after pill
    bx = 3.5; by = 3.6; bw = 2.8; bh = 1.0
    add_rect(sl, bx,        by, bw,   bh,   fill=RGBColor(0x4A,0x10,0x10), line_color=RGBColor(0x60,0x20,0x20), line_width=Pt(0.5))
    add_rect(sl, bx+bw+0.3, by, bw,   bh,   fill=RGBColor(0x0A,0x30,0x18), line_color=RGBColor(0x10,0x50,0x28), line_width=Pt(0.5))
    add_textbox(sl, "BEFORE",      bx,        by+0.06, bw,   0.25, size=9,  bold=True, color=RGBColor(0xF8,0x71,0x71), align=PP_ALIGN.CENTER)
    add_textbox(sl, "5 manual steps\n10 min per building", bx, by+0.32, bw, 0.55, size=13, color=RGBColor(0xE5,0xE7,0xEB), align=PP_ALIGN.CENTER)
    add_textbox(sl, "→",           bx+bw+0.0, by+0.3,  0.35, 0.4,  size=22, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(sl, "AFTER",       bx+bw+0.3, by+0.06, bw,   0.25, size=9,  bold=True, color=RGBColor(0x6E,0xE7,0xB7), align=PP_ALIGN.CENTER)
    add_textbox(sl, "1 postcode\n<30 s automated",         bx+bw+0.3, by+0.32, bw, 0.55, size=13, color=RGBColor(0xE5,0xE7,0xEB), align=PP_ALIGN.CENTER)

    # Stats row
    stats = [("5", "CIRCUS countries"), ("0", "Manual steps (IE)"), ("63", "Tests passing")]
    for i, (num, lbl) in enumerate(stats):
        sx = 3.8 + i * 2.0
        add_textbox(sl, num, sx, 4.85, 1.6, 0.55, size=32, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_textbox(sl, lbl, sx, 5.38, 1.6, 0.3,  size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # Footer
    add_rect(sl, 0, 7.18, 13.33, 0.32, fill=RGBColor(0x1a,0x2a,0x3a))
    add_textbox(sl, "BER Automation Tool · CIRCUS Project",
                0.22, 7.19, 7.0, 0.28, size=9, color=RGBColor(0x6B,0x72,0x80))
    add_textbox(sl, "Avinash Nagarajan · MTU · March 2026",
                6.0,  7.19, 7.0, 0.28, size=9, color=RGBColor(0x6B,0x72,0x80),
                align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Old Way
# ══════════════════════════════════════════════════════════════════════════════

def slide2(prs):
    sl = blank_slide(prs)
    header(sl, "The Old Way — Five Manual Steps Per Building", "01")
    footer(sl, "The Problem", "Slide 2 of 5")

    steps = [
        ("🗺", "Open Google Maps", "Find the address manually"),
        ("📐", "Measure footprint",  "Click-drag ruler on satellite"),
        ("🧍", "Street View survey", "Guess era, type, storeys, heating"),
        ("📊", "Enter Excel",        "Type all values into spreadsheet"),
        ("📋", "Read the answer",    "Copy result from output cell"),
    ]

    xs = [0.18, 2.72, 5.26, 7.80, 10.34]
    bw = 2.38; bh = 1.45

    for i, (icon, title, sub) in enumerate(steps):
        x = xs[i]
        add_rect(sl, x, 0.85, bw, bh, fill=RED_LIGHT, line_color=RED_BORDER, line_width=Pt(0.75))
        add_textbox(sl, icon,  x,      0.88, bw, 0.4,  size=20, align=PP_ALIGN.CENTER)
        add_textbox(sl, title, x,      1.28, bw, 0.36, size=11, bold=True,  color=RED_DARK,  align=PP_ALIGN.CENTER)
        add_textbox(sl, sub,   x,      1.62, bw, 0.60, size=10, color=GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(sl, "→", x+bw-0.05, 1.25, 0.25, 0.4, size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # Time bar
    add_rect(sl, 0.18, 2.42, 12.97, 0.38, fill=RED_LIGHT, line_color=RED_BORDER, line_width=Pt(0.75))
    add_textbox(sl, "⏱   5–10 minutes per building · 200 buildings = up to 33 hours of manual work",
                0.18, 2.43, 12.97, 0.35, size=11, bold=True, color=RED_DARK, align=PP_ALIGN.CENTER)

    # Problem boxes
    problems = [
        ("⚠  Inconsistency",   "Two assessors, same Street View,\ndifferent era estimates — every time."),
        ("⚠  Single viewpoint", "Oil tanks, heat pumps, shared walls\n— invisible from one Street View angle."),
        ("⚠  No scale",         "The CIRCUS project needs thousands\nof assessments across 5 countries."),
    ]
    pw = 4.0
    for i, (title, body) in enumerate(problems):
        px = 0.18 + i * (pw + 0.22)
        add_rect(sl, px, 2.92, pw, 1.1, fill=AMBER_LIGHT, line_color=RGBColor(0xFE,0xD7,0xAA), line_width=Pt(0.75))
        add_textbox(sl, title, px+0.12, 2.96, pw-0.2, 0.3,  size=11, bold=True, color=AMBER_DARK)
        add_textbox(sl, body,  px+0.12, 3.28, pw-0.2, 0.68, size=10, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Old vs New
# ══════════════════════════════════════════════════════════════════════════════

def slide3(prs):
    sl = blank_slide(prs)
    header(sl, "Old Way vs New Way — What Changed", "02")
    footer(sl, "Before vs After", "Slide 3 of 5")

    # Column header bars
    add_rect(sl, 0.15, 0.78, 6.4,  0.42, fill=RED)
    add_textbox(sl, "BEFORE — Excel Manual Workflow",
                0.15, 0.80, 6.4, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 6.77, 0.78, 6.4, 0.42, fill=GREEN_MID)
    add_textbox(sl, "AFTER — BER Automation Tool",
                6.77, 0.80, 6.4, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    old_steps = [
        ("Step 1 — Find address on Google Maps",        "Open browser, search address, navigate to property"),
        ("Step 2 — Measure footprint manually",          "Use satellite ruler, click around roof edge, note L × W"),
        ("Step 3 — Street View survey (one angle)",      "Visually estimate era, building type, storeys, heating"),
        ("Step 4 — Type parameters into Excel",          "Switch to spreadsheet, enter all values, risk transcription errors"),
        ("Step 5 — Read BER band from output cell",      "Copy result. No CO₂, no heat loss breakdown, no retrofit model."),
    ]
    new_steps = [
        ("Step 1 — Enter postcode, click Run",                "Eircode geocoded automatically via Google Maps API"),
        ("Step 2 — Satellite footprint extracted automatically","OpenCV + Claude Vision — two independent methods, reconciled"),
        ("Step 3 — AI surveys 4 angles simultaneously",       "Front · Right · Rear · Left → Claude cross-references in 1 call"),
        ("Step 4 — HWB calculation runs automatically",       "All parameters populated from pipeline — no manual data entry"),
        ("Step 5 — Full result displayed",                    "BER band · kWh/m²/yr · CO₂ · heat loss breakdown · retrofit model"),
    ]

    row_h = 0.94
    for i, ((ot, os_), (nt, ns_)) in enumerate(zip(old_steps, new_steps)):
        y = 1.28 + i * (row_h + 0.06)
        # OLD
        add_rect(sl, 0.15, y, 6.4, row_h, fill=RED_LIGHT, line_color=RED_BORDER, line_width=Pt(0.5))
        add_textbox(sl, ot, 0.28, y+0.08, 6.14, 0.36, size=10, bold=True,  color=RED_DARK)
        add_textbox(sl, os_, 0.28, y+0.44, 6.14, 0.44, size=10, color=DARK)
        # NEW
        add_rect(sl, 6.77, y, 6.4, row_h, fill=GREEN_LIGHT, line_color=GREEN_BORDER, line_width=Pt(0.5))
        add_textbox(sl, nt, 6.90, y+0.08, 6.14, 0.36, size=10, bold=True,  color=GREEN_DARK)
        add_textbox(sl, ns_, 6.90, y+0.44, 6.14, 0.44, size=10, color=DARK)

    # Summary bars
    sy = 1.28 + 5 * (row_h + 0.06)
    add_rect(sl, 0.15, sy, 6.4,  0.72, fill=RGBColor(0x99,0x1B,0x1B))
    add_textbox(sl, "5–10 min · 1 building · 1 assessor",       0.15, sy+0.04, 6.4, 0.3,  size=11, bold=True, color=WHITE,                       align=PP_ALIGN.CENTER)
    add_textbox(sl, "No consistency · single angle · Ireland only", 0.15, sy+0.34, 6.4, 0.3, size=10, color=RGBColor(0xFE,0xCA,0xCA), align=PP_ALIGN.CENTER)

    add_rect(sl, 6.77, sy, 6.4,  0.72, fill=GREEN_DARK)
    add_textbox(sl, "<30 s · any address · fully repeatable",   6.77, sy+0.04, 6.4, 0.3,  size=11, bold=True, color=WHITE,                        align=PP_ALIGN.CENTER)
    add_textbox(sl, "Consistent AI · 4 angles · 5 CIRCUS countries", 6.77, sy+0.34, 6.4, 0.3, size=10, color=RGBColor(0x6E,0xE7,0xB7), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What the New Tool Adds
# ══════════════════════════════════════════════════════════════════════════════

def slide4(prs):
    sl = blank_slide(prs)
    header(sl, "What the New Tool Adds Over the Excel", "03")
    footer(sl, "What's New", "Slide 4 of 5")

    cards = [
        (BLUE,   BLUE_LIGHT,  BLUE_DARK,  "🤖", "4-Angle AI Classification",
         "The Excel relied on a human guessing from one Street View angle. The tool sends Front, Right, Rear, and Left views to Claude in a single request — spotting oil tanks in the back garden, heat pump units on the side, and shared party walls.",
         "Excel: 1 angle, manual guess  →  Tool: 4 angles, AI classified"),

        (GREEN_MID, GREEN_LIGHT, GREEN_DARK, "📐", "Automatic Footprint Extraction",
         "No more clicking around a rooftop with the Google Maps ruler. Two independent automated methods — OpenCV computer vision and Claude Vision — run and reconcile. Terrace correction divides shared row dimensions by unit count.",
         "Excel: manual clicks  →  Tool: OpenCV + Claude, reconciled"),

        (PURPLE,  PURPLE_LIGHT, PURPLE_DARK, "🔧", "Retrofit Modelling",
         "The Excel produced a single BER band for the current state. The tool adds a full what-if retrofit analysis — wall insulation, roof insulation, window replacement, heating system upgrade — and shows before/after in one click.",
         "Excel: current state only  →  Tool: before + after retrofit"),

        (AMBER,   AMBER_LIGHT, AMBER_DARK,  "🌍", "5-Country Climate Data",
         "All five CIRCUS partner countries have their own HDD, solar irradiance, and heating days baked in. Germany and Ireland produce near-identical results (HDD 2,157 vs 2,149). France is noticeably milder at 1,462 HDD.",
         "Excel: manual selection  →  Tool: auto from geocoder (IE) or dropdown"),

        (RED,     RED_LIGHT,  RED_DARK,    "📊", "Full Energy Profile",
         "The Excel output was a BER band and a kWh/m² number. The tool adds annual CO₂ emissions, full heat loss breakdown (transmission vs ventilation vs solar gains), and an interactive BER scale with the property's band highlighted.",
         "Excel: BER band + kWh/m²  →  Tool: band + CO₂ + breakdown + retrofit"),

        (SLATE,   SLATE_LIGHT, RGBColor(0x1E,0x29,0x3B), "🛡", "Confidence-Gated Fallbacks",
         "When Claude can't see the building clearly — hedge in the way, no Street View coverage — it reports confidence below 0.4 and the pipeline uses safe defaults. The Excel had no such safety net; a wrong guess went straight through.",
         "Excel: wrong guess → wrong result  →  Tool: low confidence → safe defaults"),
    ]

    cw = 4.1; ch = 2.06
    xs = [0.18, 4.56, 8.94]
    ys = [0.82, 3.06]

    for idx, (border, bg, title_col, icon, title, body, pill) in enumerate(cards):
        col = idx % 3
        row = idx // 3
        x = xs[col]; y = ys[row]

        add_rect(sl, x, y, cw, ch, fill=bg, line_color=None)
        # Top border accent (simulate border-top via thin rect)
        add_rect(sl, x, y, cw, 0.07, fill=border)

        add_textbox(sl, icon,  x+0.14, y+0.10, 0.5, 0.38, size=18)
        add_textbox(sl, title, x+0.14, y+0.50, cw-0.25, 0.34, size=11, bold=True, color=title_col)
        add_textbox(sl, body,  x+0.14, y+0.84, cw-0.25, 0.80, size=9.5, color=DARK)
        add_rect(sl, x+0.10, y+ch-0.34, cw-0.20, 0.28, fill=RGBColor(0xE0,0xE8,0xFF) if border==BLUE else
                 GREEN_BORDER if border==GREEN_MID else
                 RGBColor(0xE9,0xD5,0xFF) if border==PURPLE else
                 RGBColor(0xFD,0xE6,0x8A) if border==AMBER else
                 RGBColor(0xFE,0xCD,0xD3) if border==RED else
                 RGBColor(0xCB,0xD5,0xE0))
        add_textbox(sl, pill, x+0.14, y+ch-0.32, cw-0.25, 0.26, size=8.5, bold=True, color=title_col)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What-If / Community Engagement
# ══════════════════════════════════════════════════════════════════════════════

def slide5(prs):
    sl = blank_slide(prs)
    header(sl, "The What-If Engine — Engaging Communities as Energy Citizens", "04")
    footer(sl, "Community Engagement", "Slide 5 of 5")

    # ── Before retrofit box ───────────────────────────────────────────────────
    bx = 0.18; by = 0.85; bw = 3.8; bh = 2.55
    add_rect(sl, bx, by, bw, bh, fill=RED_LIGHT, line_color=RED_BORDER, line_width=Pt(0.75))
    add_textbox(sl, "BEFORE RETROFIT", bx, by+0.08, bw, 0.26,
                size=10, bold=True, color=RED_DARK, align=PP_ALIGN.CENTER)
    # Band badge
    add_rect(sl, bx+0.18, by+0.42, 1.0, 0.9, fill=ORANGE)
    add_textbox(sl, "D2", bx+0.18, by+0.52, 1.0, 0.62,
                size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "band", bx+0.18, by+1.10, 1.0, 0.22,
                size=9, color=WHITE, align=PP_ALIGN.CENTER)
    # Stats
    add_textbox(sl, "286 kWh/m²/yr",     bx+1.35, by+0.46, 2.2, 0.3,  size=12, bold=True, color=DARK)
    add_textbox(sl, "4,100 kg CO₂/yr",   bx+1.35, by+0.76, 2.2, 0.28, size=10, color=GRAY)
    add_textbox(sl, "Pre-1980 · oil boiler", bx+1.35, by+1.02, 2.2, 0.28, size=10, color=GRAY)
    add_textbox(sl, "Solid walls · single glaze", bx+1.35, by+1.28, 2.2, 0.28, size=10, color=GRAY)
    # Cost callout
    add_rect(sl, bx+0.14, by+1.88, bw-0.28, 0.40, fill=RGBColor(0xFE,0xE2,0xE2))
    add_textbox(sl, "≈ €2,400/yr heating cost", bx+0.14, by+1.90, bw-0.28, 0.20,
                size=9, bold=True, color=RED_DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Among worst 20% of Irish stock", bx+0.14, by+2.10, bw-0.28, 0.20,
                size=9, color=RED_DARK, align=PP_ALIGN.CENTER)

    # ── Arrow + interventions ─────────────────────────────────────────────────
    mx = 4.2; my = 0.85; mw = 2.9; mh = 2.55
    add_textbox(sl, "→", mx-0.05, by+0.95, 0.5, 0.6,
                size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_rect(sl, mx+0.5, my, mw, mh, fill=PURPLE_LIGHT, line_color=RGBColor(0xDD,0xD6,0xFE), line_width=Pt(0.75))
    add_textbox(sl, "MEASURES APPLIED", mx+0.5, my+0.08, mw, 0.26,
                size=10, bold=True, color=PURPLE_DARK, align=PP_ALIGN.CENTER)
    measures = ["🧱  12 cm wall insulation",
                "🏠  20 cm roof insulation",
                "🪟  New windows  U=1.0",
                "♨️  Air-source heat pump",
                "     (replaces oil boiler)"]
    for i, m in enumerate(measures):
        add_textbox(sl, m, mx+0.65, my+0.46+i*0.38, mw-0.3, 0.35,
                    size=10, color=DARK if i < 4 else GRAY, italic=(i==4))

    # ── After retrofit box ────────────────────────────────────────────────────
    ax = 7.55; ay = 0.85; aw = 3.8; ah = 2.55
    add_rect(sl, ax, ay, aw, ah, fill=GREEN_LIGHT, line_color=GREEN_BORDER, line_width=Pt(0.75))
    add_textbox(sl, "→", ax-0.6, ay+0.95, 0.6, 0.6,
                size=30, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "AFTER RETROFIT", ax, ay+0.08, aw, 0.26,
                size=10, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    # Band badge
    add_rect(sl, ax+0.18, ay+0.42, 1.0, 0.9, fill=GREEN_MID)
    add_textbox(sl, "B1", ax+0.18, ay+0.52, 1.0, 0.62,
                size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "band", ax+0.18, ay+1.10, 1.0, 0.22,
                size=9, color=WHITE, align=PP_ALIGN.CENTER)
    # Stats
    add_textbox(sl, "118 kWh/m²/yr",         ax+1.35, ay+0.46, 2.2, 0.3,  size=12, bold=True, color=DARK)
    add_textbox(sl, "1,200 kg CO₂/yr  (−71%)", ax+1.35, ay+0.76, 2.2, 0.28, size=10, color=GRAY)
    add_textbox(sl, "5-band improvement",      ax+1.35, ay+1.02, 2.2, 0.28, size=10, bold=True, color=GREEN_MID)
    add_textbox(sl, "Energy use halved",        ax+1.35, ay+1.28, 2.2, 0.28, size=10, bold=True, color=GREEN_MID)
    # Saving callout
    add_rect(sl, ax+0.14, ay+1.88, aw-0.28, 0.40, fill=RGBColor(0xDC,0xFC,0xE7))
    add_textbox(sl, "≈ €1,500/yr saved on bills", ax+0.14, ay+1.90, aw-0.28, 0.20,
                size=9, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Carbon footprint cut by two-thirds", ax+0.14, ay+2.10, aw-0.28, 0.20,
                size=9, color=GREEN_DARK, align=PP_ALIGN.CENTER)

    # ── Three community cards ─────────────────────────────────────────────────
    community_cards = [
        (BLUE,     BLUE_LIGHT,   BLUE_DARK,   "🏘️", "From Data to Decision",
         "A BER band alone is abstract. The what-if model converts it into a concrete question residents understand: \"What would it take to halve my heating bill?\" The tool answers that for any specific home — without an assessor visit."),

        (GREEN_MID, GREEN_LIGHT, GREEN_DARK,  "⚡", "Enabling Energy Citizens",
         "Energy citizenship requires informed participation. When communities can see their own street — every house rated and modelled — they can identify vulnerable neighbours, plan area-wide retrofit schemes, and present evidence to local authorities."),

        (PURPLE,   PURPLE_LIGHT, PURPLE_DARK, "📍", "CIRCUS Scale: Estate-Level Screening",
         "Running the what-if model on every house in a 200-home estate takes minutes. The output tells planners which interventions deliver the biggest gain per euro — wall insulation or heating upgrades, and in what order."),
    ]

    cy = 3.58; ch = 3.38; cw = 4.1
    cxs = [0.18, 4.61, 9.04]

    for i, (border, bg, title_col, icon, title, body) in enumerate(community_cards):
        cx = cxs[i]
        add_rect(sl, cx, cy, cw, ch, fill=bg, line_color=None)
        add_rect(sl, cx, cy, cw, 0.07, fill=border)   # top accent
        add_textbox(sl, icon,  cx+0.14, cy+0.14, 0.5,      0.42, size=18)
        add_textbox(sl, title, cx+0.14, cy+0.60, cw-0.25,  0.36, size=12, bold=True,  color=title_col)
        add_textbox(sl, body,  cx+0.14, cy+1.00, cw-0.25,  2.22, size=10, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    prs = new_prs()
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)

    out = os.path.join(os.path.dirname(__file__), "BER_Automation_CIRCUS.pptx")
    prs.save(out)
    print(f"Saved: {out}")
