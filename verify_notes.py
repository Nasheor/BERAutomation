from pptx import Presentation
prs = Presentation(r"C:\Users\nagar\PyCharmProjects\BERAutomation\docs\environ_2026_presentation.pptx")
for i, slide in enumerate(prs.slides, 1):
    text = slide.notes_slide.notes_text_frame.text.strip()
    preview = text[:65].replace("\n", " ")
    print(f"Slide {i:2d}: {len(text):4d} chars  |  {preview}...")
